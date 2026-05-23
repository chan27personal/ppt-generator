# -*- coding: utf-8 -*-
"""pptlib — 표준 PPT 제작 헬퍼 (python-pptx). 범용(크로스플랫폼).
A4 세로/가로·16:9, 모던(로열블루)/네이비 팔레트, 라벨형 불릿, 아이소메트릭 블록.
- 폰트: 이름만 주면 OS 폰트 폴더에서 파일 자동 탐색. 없으면 줄바꿈을 휴리스틱으로 추정(렌더는 PPT 폰트 치환).
- 한글은 latin + East Asian(a:ea) 모두 지정.

예:
    from pptlib import Deck, Modern, Navy, MSO_SHAPE, PP_ALIGN, MSO_ANCHOR
    d = Deck(page="A4P", theme=Modern, font="사천항공")   # font 생략 시 OS 기본
    s = d.page()
    d.text(0.5, 0.5, 5, 0.4, "제목", 22, d.t.INK, bold=True)
    d.save("out.pptx")
"""
import os, glob, platform
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import ImageFont


def rgb(h): return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Modern:   # 로열블루 (기본 권장 — 모던 제안서/덱)
    PRI = rgb("3B5BF0"); PRI_D = rgb("2740C4"); PRI_L = rgb("6E88F7")
    INK = rgb("1B2233"); SUB = rgb("5A6373")
    TINT = rgb("EEF1FE"); TINT2 = rgb("E0E6FC"); LINE = rgb("E1E7F6"); WHITE = rgb("FFFFFF")


class Navy:     # 네이비/스카이 (정형 공공·기업 문서)
    PRI = rgb("2E78CB"); PRI_D = rgb("1B3663"); PRI_L = rgb("4A90D9")
    INK = rgb("1B3663"); SUB = rgb("555E6B")
    TINT = rgb("EAF1FB"); TINT2 = rgb("DDE9F9"); LINE = rgb("C7D7EC"); WHITE = rgb("FFFFFF")


PAGE = {"A4P": (8.27, 11.69), "A4L": (11.69, 8.27), "16:9": (13.333, 7.5)}

# ---------- 폰트 자동 탐색 (범용) ----------
def _font_dirs():
    s = platform.system(); d = []
    if s == "Windows":
        d += [os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts"),
              os.path.join(os.environ.get("LOCALAPPDATA", ""), "Microsoft", "Windows", "Fonts")]
    elif s == "Darwin":
        d += ["/System/Library/Fonts", "/System/Library/Fonts/Supplemental",
              "/Library/Fonts", os.path.expanduser("~/Library/Fonts")]
    else:
        d += ["/usr/share/fonts", "/usr/local/share/fonts",
              os.path.expanduser("~/.fonts"), os.path.expanduser("~/.local/share/fonts")]
    return [x for x in d if x and os.path.isdir(x)]

_KNOWN = {   # 표시이름(소문자, 공백제거) -> 후보 파일명
    "사천항공": ["SacheonHangGong-Regular.ttf", "SacheonHangGong.ttf"],
    "맑은고딕": ["malgun.ttf"], "malgungothic": ["malgun.ttf"],
    "나눔고딕": ["NanumGothic.ttf"], "nanumgothic": ["NanumGothic.ttf"],
    "pretendard": ["Pretendard-Regular.otf", "Pretendard-Regular.ttf"],
    "notosanskr": ["NotoSansKR-Regular.otf", "NotoSansKR-Regular.ttf", "NotoSansCJKkr-Regular.otf"],
    "notosanscjkkr": ["NotoSansCJKkr-Regular.otf", "NotoSansCJK-Regular.ttc"],
    "applesdgothicneo": ["AppleSDGothicNeo.ttc"],
    "arial": ["arial.ttf", "Arial.ttf"], "calibri": ["calibri.ttf"],
    "helvetica": ["Helvetica.ttc"],
}

def _find_font_file(name):
    if not name:
        return None
    key = name.lower().replace(" ", "")
    cands = _KNOWN.get(key, [])
    dirs = _font_dirs()
    for fn in cands:                      # 알려진 파일명 우선
        for d in dirs:
            for p in glob.glob(os.path.join(d, "**", fn), recursive=True):
                return p
            direct = os.path.join(d, fn)
            if os.path.exists(direct):
                return direct
    tok = key.replace("-", "")            # 파일명 토큰 매칭(차선)
    if tok:
        for d in dirs:
            for ext in ("*.ttf", "*.ttc", "*.otf"):
                for p in glob.glob(os.path.join(d, "**", ext), recursive=True):
                    base = os.path.basename(p).lower().replace(" ", "").replace("-", "")
                    if tok in base:
                        return p
    return None

def _default_font():
    if _find_font_file("사천항공"):       # 이 환경의 선호 폰트가 있으면 사용
        return "사천항공"
    s = platform.system()
    return {"Windows": "Malgun Gothic", "Darwin": "Apple SD Gothic Neo"}.get(s, "Noto Sans CJK KR")


class Deck:
    def __init__(self, page="A4P", font=None, font_file=None, theme=Modern):
        w, h = PAGE[page]
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(w * 914400)); self.prs.slide_height = Emu(int(h * 914400))
        self.PW, self.PH = w, h
        self.font = font or _default_font()
        self.t, self._fc = theme, {}
        # 측정용 폰트 파일: 지정 > 본문폰트 탐색 > 일반 폴백
        self._mfile = (font_file or _find_font_file(self.font)
                       or _find_font_file("Malgun Gothic") or _find_font_file("Arial")
                       or _find_font_file("Noto Sans CJK KR"))
        self.slide = None

    def page(self):
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); return self.slide

    # ---------- 텍스트 폭/줄바꿈 (폰트 없으면 휴리스틱) ----------
    def _pil(self, sz):
        if not self._mfile:
            return None
        px = max(6, int(round(sz * 96 / 72)))
        if px not in self._fc:
            try:
                self._fc[px] = ImageFont.truetype(self._mfile, px)
            except Exception:
                self._fc[px] = None
        return self._fc[px]

    def _w_px(self, s, sz):
        f = self._pil(sz)
        if f is not None:
            return f.getlength(s)
        em = sz * 96 / 72                 # 휴리스틱: CJK=1em, 라틴=0.52em, 공백=0.30em
        tot = 0.0
        for ch in s:
            o = ord(ch)
            tot += em * (1.0 if (0x1100 <= o <= 0x11FF or 0x2E80 <= o <= 0x9FFF or 0xAC00 <= o <= 0xD7A3 or 0xF900 <= o <= 0xFAFF)
                         else (0.30 if ch == " " else 0.52))
        return tot

    def wrap(self, s, sz, w_in):
        mx = w_in * 96 * 0.95; out, cur = [], ""
        for word in s.split(" "):
            t = word if cur == "" else cur + " " + word
            if cur == "" or self._w_px(t, sz) <= mx: cur = t
            else: out.append(cur); cur = word
            while self._w_px(cur, sz) > mx and len(cur) > 1:
                p = ""
                for ch in cur:
                    if p and self._w_px(p + ch, sz) > mx: break
                    p += ch
                out.append(p); cur = cur[len(p):]
        if cur != "": out.append(cur)
        return out or [""]

    def nlines(self, s, sz, w_in): return len(self.wrap(s, sz, w_in))

    # ---------- 기본 도형/텍스트 ----------
    def _ea(self, r):
        rPr = r._r.get_or_add_rPr(); el = rPr.find(qn('a:ea'))
        if el is None: el = rPr.makeelement(qn('a:ea'), {}); rPr.append(el)
        el.set('typeface', self.font)

    def shp(self, kind, x, y, w, h, fill=None, line=None, lw=1.0, radius=None):
        s = self.slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None: s.fill.background()
        else: s.fill.solid(); s.fill.fore_color.rgb = fill
        if line is None: s.line.fill.background()
        else: s.line.color.rgb = line; s.line.width = Pt(lw)
        s.shadow.inherit = False
        if radius is not None:
            try: s.adjustments[0] = radius
            except Exception: pass
        return s

    def text(self, x, y, w, h, content, sz, color, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, ls=None):
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.03); tf.margin_top = tf.margin_bottom = Emu(0)
        for i, ln in enumerate(content.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if ls: p.line_spacing = Pt(ls)
            p.space_before = Pt(0); p.space_after = Pt(0)
            r = p.add_run(); r.text = ln
            r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = self.font; self._ea(r)
        return tb

    def rich(self, x, y, w, h, runs, sz, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, ls=None):
        """runs = [(text, color, bold), ...] — 한 줄에 여러 색/굵기."""
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.02); tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = align
        if ls: p.line_spacing = Pt(ls)
        for txt, color, bold in runs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = self.font; self._ea(r)
        return tb

    def pic(self, png, cx, cy, side):
        """중심 좌표 기준 정사각 이미지(아이콘)."""
        self.slide.shapes.add_picture(png, Inches(cx - side / 2), Inches(cy - side / 2),
                                      Inches(side), Inches(side))

    def chevron(self, cx, cy, size, color=None, lw=2.2):
        """폰트 글리프 대신 선 2개로 그린 '›'."""
        color = color or self.t.PRI; h = size / 2
        for (ax, ay), (bx, by) in [((cx - h * 0.5, cy - h), (cx + h * 0.5, cy)),
                                    ((cx + h * 0.5, cy), (cx - h * 0.5, cy + h))]:
            cn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                                 Inches(ax), Inches(ay), Inches(bx), Inches(by))
            cn.line.color.rgb = color; cn.line.width = Pt(lw); cn.shadow.inherit = False

    # ---------- 라벨형 불릿 ----------
    def rbullet(self, x, y, w, label, desc, sz, ls=None, lab_c=None, desc_c=None,
                mk=None, ind=0.14, gap=0.03, sep=" : "):
        lab_c = lab_c or self.t.PRI; desc_c = desc_c or self.t.INK; mk = mk or self.t.PRI
        if ls is None: ls = sz * 1.34
        full = label + sep + desc; tw = w - ind; lines = self.wrap(full, sz, tw)
        m = 0.05
        self.shp(MSO_SHAPE.OVAL, x, y + (ls / 72) / 2 - m / 2, m, m, fill=mk, line=None)
        h = len(lines) * ls / 72
        tb = self.slide.shapes.add_textbox(Inches(x + ind), Inches(y), Inches(tw), Inches(h + 0.02))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        ll = len(label)
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT; p.line_spacing = Pt(ls)
            p.space_before = Pt(0); p.space_after = Pt(0)
            if i == 0:
                r1 = p.add_run(); r1.text = label
                r1.font.size = Pt(sz); r1.font.bold = True; r1.font.color.rgb = lab_c
                r1.font.name = self.font; self._ea(r1)
                rest = ln[ll:]; r2 = p.add_run(); r2.text = rest if rest != "" else " "
                r2.font.size = Pt(sz); r2.font.color.rgb = desc_c
                r2.font.name = self.font; self._ea(r2)
            else:
                r = p.add_run(); r.text = ln
                r.font.size = Pt(sz); r.font.color.rgb = desc_c
                r.font.name = self.font; self._ea(r)
        return y + h + gap

    def rbh(self, w, label, desc, sz, ls=None, ind=0.14, gap=0.03, sep=" : "):
        if ls is None: ls = sz * 1.34
        return self.nlines(label + sep + desc, sz, w - ind) * ls / 72 + gap

    # ---------- 아이소메트릭 블록 ----------
    def iso_box(self, cx, cy, w, t, top=None, left=None, right=None):
        top = top or self.t.PRI_L; left = left or self.t.PRI; right = right or self.t.PRI_D

        def emu(v): return int(round(v * 914400))

        def poly(pts, fill):
            pe = [(emu(a), emu(b)) for a, b in pts]
            ff = self.slide.shapes.build_freeform(pe[0][0], pe[0][1], scale=1.0)
            ff.add_line_segments(pe[1:], close=True)
            s = ff.convert_to_shape(); s.fill.solid(); s.fill.fore_color.rgb = fill
            s.line.color.rgb = fill; s.line.width = Pt(0.4); s.shadow.inherit = False
        half, q = w / 2, w / 4
        poly([(cx - half, cy), (cx, cy + q), (cx, cy + q + t), (cx - half, cy + t)], left)
        poly([(cx, cy + q), (cx + half, cy), (cx + half, cy + t), (cx, cy + q + t)], right)
        poly([(cx, cy - q), (cx + half, cy), (cx, cy + q), (cx - half, cy)], top)

    def iso_stack(self, cx, cy0, w, t, n=4, gap=0.40):
        for i in range(n): self.iso_box(cx, cy0 + i * gap, w, t)

    def eyebrow(self, x, y, text, w, sz=12.5):
        self.shp(MSO_SHAPE.RECTANGLE, x, y + 0.02, 0.20, 0.20, fill=self.t.PRI, line=None)
        self.text(x + 0.30, y - 0.02, w - 0.3, 0.28, text, sz, self.t.INK, bold=True, align=PP_ALIGN.LEFT)
        return y + 0.36

    def save(self, path):
        self.prs.save(path); return path
